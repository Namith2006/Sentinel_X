from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def apply_dark_theme(slide, prs):
    # Set dark slate background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 25, 35) # Deep slate/navy
    
    # Add an emerald green accent bar at the top
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Inches(0.15)
    shape = slide.shapes.add_shape(
        1, left, top, width, height # 1 is rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(16, 185, 129) # Emerald Green
    shape.line.color.rgb = RGBColor(16, 185, 129)

def main():
    prs = Presentation()
    
    # Slides content matching the MES College rubric + New additions
    slides_data = [
        {
            "title": "Sentinel X AI Engine",
            "content": [
                "Project Work Introductory Seminar",
                "Department of Computer Applications",
                "Team: Balaji R, Namith M R, Jagadeesh N"
            ]
        },
        {
            "title": "Agenda of your presentation",
            "content": [
                "Introduction to Sentinel X",
                "Existing Systems & Literature Survey",
                "Limitations of Existing Systems",
                "Proposed System, Objective & Methodology",
                "Hardware & Software Requirements",
                "Core Modules Identified",
                "Applicability in the Present Scenario",
                "Team Roles & Contributions",
                "Future Scope and Advancements",
                "Visualizations of End Product"
            ]
        },
        {
            "title": "Introduction",
            "content": [
                "Sentinel X is a hybrid cybersecurity AI dashboard.",
                "Designed to detect zero-day phishing links and synthetic media (deepfakes).",
                "Integrates machine learning with deterministic heuristic rules.",
                "Provides automated, natural-language incident mitigation for end users."
            ]
        },
        {
            "title": "Existing System",
            "content": [
                "Standard cybersecurity scanners rely on basic lexical URL patterns.",
                "Traditional image forensics rely purely on raw pixel metadata (EXIF).",
                "Existing threat dashboards only alert users but offer no active recovery steps."
            ]
        },
        {
            "title": "Brief on the literature survey carried out",
            "content": [
                "Research concludes CNNs (Convolutional Neural Networks) effectively detect raw GAN artifacts.",
                "Studies show deepfake detection accuracy severely drops post-JPEG compression.",
                "Research indicates lexical URL scanners consistently miss third-party adware distributors.",
                "Conclusion: AI alone is insufficient; hybrid heuristic models are required for edge cases."
            ]
        },
        {
            "title": "Limitations of Existing system",
            "content": [
                "Disadvantage 1: Standard AI fails on highly compressed media (WhatsApp/Telegram).",
                "Disadvantage 2: Scanners miss hidden zero-day malware on legal domains (Softonic, APK mirrors).",
                "Disadvantage 3: High false-positive rates on synthetic-looking system UI screenshots.",
                "Disadvantage 4: Users are left stranded without technical guidance after a threat is found."
            ]
        },
        {
            "title": "Proposed System, Its objective and Methodology",
            "content": [
                "Objective: Build a hybrid AI that overcomes edge-case bypasses and guides users.",
                "Methodology & How it Overcomes Limitations:",
                "1. Bypasses WhatsApp compression using Error Level Analysis (ELA) and Laplacian variance.",
                "2. Overrides AI URL false-negatives with aggressive heuristic APK/malware keyword engines.",
                "3. Implements screenshot logic to prevent false positives on flat digital UI graphics.",
                "4. Triggers local Llama-3 LLM to generate instant, 3-step JSON mitigation action plans."
            ]
        },
        {
            "title": "Hardware and Software requirements",
            "content": [
                "Hardware Requirements:",
                "- High-performance processor (Core i7 / Ryzen 7 or equivalent).",
                "- Dedicated GPU for local tensor computation (e.g., RTX 40-series / MSI Crosshair 16 HX).",
                "- Minimum 16GB RAM for local LLM inference.",
                "Software Requirements:",
                "- Python 3.11+",
                "- FastAPI, Uvicorn, PyTorch, OpenCV",
                "- Ollama (Local Llama-3 Runtime)",
                "- Frontend: HTML5, Tailwind CSS, JavaScript (Fetch API)"
            ]
        },
        {
            "title": "Modules identified",
            "content": [
                "Phishing & URL Scanner Engine (phishing_engine.py)",
                "Deepfake & Image Forensics Engine (deepfake_engine.py)",
                "AI Mitigation Chatbot Expert (llm_expert.py)",
                "FastAPI REST Gateway & Crypto Ledger (main.py)",
                "WebAuthn & Interactive Dashboard (index.html)"
            ]
        },
        {
            "title": "Applicability of the project in the present scenario",
            "content": [
                "Protects everyday users from modern social engineering and deepfake scams.",
                "Intercepts malicious third-party APK downloads common on messaging apps.",
                "Provides enterprise-grade threat intelligence in an accessible, local environment.",
                "Requires zero internet connection for the AI engine, ensuring total data privacy."
            ]
        },
        {
            "title": "Team Roles & Contributions",
            "content": [
                "Balaji R (Frontend & Architecture): Designed dark-mode UI, implemented WebAuthn Passkeys, and built API JSON data piping.",
                "Namith M R (AI & ML Core): Built PyTorch vision model, integrated ELA for WhatsApp compression, and engineered screenshot heuristics.",
                "Jagadeesh N (Backend & Integration): Developed FastAPI backend, engineered the Phishing URL rules, and integrated the Llama-3 expert system."
            ]
        },
        {
            "title": "Future Scope and Advancements",
            "content": [
                "Development of a web browser extension for real-time phishing and deepfake interception.",
                "Expansion of machine learning models to detect deepfake audio and voice cloning.",
                "Migration to a scalable cloud architecture for enterprise API access.",
                "Integration of real-time global threat feeds to dynamically update heuristic blocklists."
            ]
        },
        {
            "title": "Visualizations of end product",
            "content": [
                "(Paste your Dashboard, Security Ledger, and AI Mitigation screenshots here)"
            ]
        }
    ]

    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    for i, slide_data in enumerate(slides_data):
        # Use title layout for the first slide, bullet layout for the rest
        if i == 0:
            slide = prs.slides.add_slide(title_slide_layout)
            apply_dark_theme(slide, prs)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_data["title"]
            
            # Style Title
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(6, 182, 212) # Neon Cyan
            title.text_frame.paragraphs[0].font.bold = True
            
            # Style Subtitle
            subtitle.text = "\n".join(slide_data["content"])
            for p in subtitle.text_frame.paragraphs:
                p.font.color.rgb = RGBColor(255, 255, 255) # White
        else:
            slide = prs.slides.add_slide(bullet_slide_layout)
            apply_dark_theme(slide, prs)
            
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = slide_data["title"]
            
            # Style Title
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(6, 182, 212) # Neon Cyan
            title.text_frame.paragraphs[0].font.bold = True
            
            # Add bullet points
            tf = body.text_frame
            tf.text = slide_data["content"][0]
            
            for point in slide_data["content"][1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
            
            # Style Bullet Points
            for p in tf.paragraphs:
                p.font.color.rgb = RGBColor(230, 235, 240) # Off-White
                p.font.size = Pt(20)

    prs.save('Sentinel_X_Presentation.pptx')
    print("Successfully generated Sentinel_X_Presentation.pptx with the updated Literature Survey and Future Scope!")

if __name__ == "__main__":
    main()